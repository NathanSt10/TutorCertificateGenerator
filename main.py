import logging
from typing import Optional
import openpyxl
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import os

# Constants
CERTIFICATE_FONT_SIZE = 40
CERTIFICATE_FONT_NAME = "Century Schoolbook"
NAME_SHAPE_INDEX = 4  # Index of the shape containing the name field
EXCEL_EXTENSION = ".xlsx"
POWERPOINT_EXTENSION = ".pptx"
SKIP_SHEET_NAME = "Finished Level 3"

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)


# raw_name: Last, First
def parse_name(raw_name: str) -> str:
    """
       Parse a name from "Last, First" format to "First Last" format.

       Args:
           raw_name: Name in "Last, First" format

       Returns:
           Name in "First Last" format
    """
    parts = raw_name.split(',')
    if len(parts) == 2:
        last_name = parts[0].strip()
        first_name = parts[1].strip()
        return f"{first_name} {last_name}"
    return raw_name  # Return original name if it doesn't match the expected format


def create_certificate(template_path: str, name: str, output_path: str) -> None:
    """
        Create a certificate for the given name and save it to the specified path

        Args:
            template_path: Path to the PowerPoint template
            name: Name to insert in the certificate
            output_path: Path where to save the certificate
    """
    try:
        # Create a new presentation from the template for each certificate
        certificate = Presentation(template_path)
        # Get the first (only) slide in the template
        slide = certificate.slides[0]

        # Get the text box where the name should be inserted
        insert_name_tf = slide.shapes[NAME_SHAPE_INDEX].text_frame
        insert_name_tf.clear()
        insert_name_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = insert_name_tf.paragraphs[0]
        p.text = name

        run = p.runs[0]
        run.font.size = Pt(CERTIFICATE_FONT_SIZE)
        run.font.bold = True
        run.font.name = CERTIFICATE_FONT_NAME

        certificate.save(output_path)
        logger.info(f"Created certificate for {name}")
    except Exception as e:
        logger.error(f"Failed to create certificate for {name}: {e}")


def find_excel_file(directory: str) -> Optional[str]:
    """
    Find a single Excel file in the specified directory.

    Args:
        directory: Directory to search in

    Returns:
        Path to the Excel file if exactly one is found, None otherwise
    """
    xlsx_files = [f for f in os.listdir(directory) if f.endswith(EXCEL_EXTENSION)]

    if len(xlsx_files) == 1:
        return os.path.join(directory, xlsx_files[0])
    elif not xlsx_files:
        logger.error(f"No {EXCEL_EXTENSION} files found in {directory}")
    else:
        logger.error(f"Multiple {EXCEL_EXTENSION} files found in {directory}. Expected exactly one.")

    return None


def main():
    path = os.getcwd()
    xlsx_path = os.path.join(path, "attendance")

    xlsx_files = [f for f in os.listdir(xlsx_path) if f.endswith('.xlsx')]

    if len(xlsx_files) == 1:
        file_path = os.path.join(xlsx_path, xlsx_files[0])
        workbook = openpyxl.load_workbook(file_path)
        print(f"Opened: {file_path}")
    else:
        print("Error: Expected exactly one .xlsx file in the directory.")
        exit(1)

    for sheet in workbook:
        if sheet.title == "Finished Level 3":
            continue
        os.makedirs(sheet.title, exist_ok=True)
        directory = os.path.join(path, sheet.title)
        template_name = os.path.join(path, "templates", f"{sheet.title}{POWERPOINT_EXTENSION}")
        for row in sheet.iter_rows(min_row=2, max_col=1):
            for cell in row:
                if cell.value is not None and cell.font.b:
                    tutor_name = parse_name(cell.value)
                    create_certificate(template_name, tutor_name, directory)


if __name__ == "__main__":
    main()
